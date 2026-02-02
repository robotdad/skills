#!/usr/bin/env python3
"""
Convert HTML presentations to PowerPoint format.

Handles both Amplifier Stories HTML structure and generic HTML presentations.
Uses python-pptx for clear API and accurate positioning without browser overhead.

Usage:
    python convert.py input.html output.pptx
    python convert.py input.html output.pptx --primary-color 667eea --accent-color f093fb
"""

import argparse
import sys
from pathlib import Path

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    )


class HTMLToPPTXConverter:
    """Convert HTML presentations to PowerPoint."""

    def __init__(
        self,
        html_content: str,
        primary_color: str | None = None,
        accent_color: str | None = None,
    ):
        self.soup = BeautifulSoup(html_content, "lxml")
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self.blank_layout = self.prs.slide_layouts[6]

        # Colors - defaults match Amplifier Stories
        self.primary = hex_to_rgb(primary_color if primary_color else "0078D4")
        self.accent = hex_to_rgb(accent_color if accent_color else "50E6FF")
        self.black = RGBColor(0x00, 0x00, 0x00)
        self.white = RGBColor(0xFF, 0xFF, 0xFF)
        self.gray_70 = RGBColor(0xB3, 0xB3, 0xB3)
        self.gray_50 = RGBColor(0x80, 0x80, 0x80)
        self.dark_gray = RGBColor(0x1A, 0x1A, 0x1A)
        self.border_gray = RGBColor(0x33, 0x33, 0x33)

    def get_text(self, element: Tag | None) -> str:
        """Extract text from element, handling None."""
        return element.get_text(strip=True) if element else ""

    def set_background(self, slide, color: RGBColor | None = None):
        """Set solid background for slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color or self.black

    def add_text_box(
        self,
        slide,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int = 14,
        bold: bool = False,
        color: RGBColor | None = None,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        """Add text box with styling."""
        if not text:
            return None

        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or self.white
        p.alignment = align
        return box

    def is_centered(self, slide_div: Tag) -> bool:
        """Check if slide should be center-aligned."""
        classes = slide_div.get("class", [])
        return "center" in classes

    def process_slide(self, slide_div: Tag):
        """Process a single slide div and add to presentation."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        self.set_background(slide)

        is_centered = self.is_centered(slide_div)
        current_top = 0.6

        # Section label
        section_label = slide_div.find(class_="section-label")
        if section_label:
            if is_centered:
                current_top = 1.5
            self.add_text_box(
                slide,
                self.get_text(section_label).upper(),
                0.8,
                current_top,
                8.4,
                0.4,
                font_size=14,
                bold=True,
                color=self.primary,
            )
            current_top += 0.5

        # Headline - try Amplifier Stories classes first, then any h1/h2
        headline = slide_div.find(["h1", "h2"], class_="headline") or slide_div.find("h1") or slide_div.find("h2")
        if headline:
            text = self.get_text(headline)
            has_gradient = "big-text" in headline.get("class", [])
            color = self.accent if has_gradient else self.white
            size = 56 if headline.name == "h1" or has_gradient else 40

            if is_centered:
                current_top = max(current_top, 2.0)

            self.add_text_box(
                slide,
                text,
                0.8,
                current_top,
                8.4,
                1.5,
                font_size=size,
                bold=True,
                color=color,
                align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
            )
            current_top += 1.2 if size > 45 else 0.9

        # Medium headline
        medium_headline = slide_div.find(class_="medium-headline")
        if medium_headline and medium_headline != headline:
            self.add_text_box(
                slide,
                self.get_text(medium_headline),
                0.8,
                current_top,
                8.4,
                0.8,
                font_size=36,
                bold=True,
                color=self.white,
                align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
            )
            current_top += 0.8

        # Subhead
        subhead = slide_div.find(class_="subhead")
        if subhead:
            self.add_text_box(
                slide,
                self.get_text(subhead),
                0.8,
                current_top,
                8.4,
                1.0,
                font_size=24,
                color=self.gray_70,
                align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
            )
            current_top += 0.8

        # Cards
        card_containers = slide_div.find_all(class_=["thirds", "halves", "fourths"])
        for container in card_containers:
            cards = container.find_all(class_="card")
            if cards:
                self._add_cards(slide, cards, current_top)
                current_top += 2.0

        # Standalone cards
        standalone_cards = [
            c
            for c in slide_div.find_all(class_="card")
            if not c.find_parent(class_=["thirds", "halves", "fourths"])
        ]
        if standalone_cards:
            self._add_cards(slide, standalone_cards, current_top)
            current_top += 2.0

        # Tenet boxes
        tenets = slide_div.find_all(class_="tenet")
        if tenets:
            self._add_tenets(slide, tenets, current_top)
            current_top += len(tenets) * 0.5 + 0.5

        # Stats grid (try both Amplifier naming and generic)
        stat_grid = slide_div.find(class_="stat-grid") or slide_div.find(class_="stats-grid")
        if stat_grid:
            self._add_stats(slide, stat_grid, current_top)
            current_top += 2.5

        # Tables
        tables = slide_div.find_all("table", class_="data-table")
        for table in tables:
            self._add_table(slide, table, current_top)
            current_top += 2.5

        # Feature lists (Amplifier Stories format)
        feature_lists = slide_div.find_all(class_="feature-list")
        for fl in feature_lists:
            if not fl.find_parent(class_="versus"):
                self._add_feature_list(slide, fl, current_top)
                current_top += 1.5

        # Generic content fallback - extract plain HTML elements not already handled
        # This catches custom HTML that doesn't use Amplifier Stories classes
        self._add_generic_content(slide, slide_div, current_top, is_centered)

        # Small text at bottom
        small_text = slide_div.find(class_="small-text")
        if small_text:
            self.add_text_box(
                slide,
                self.get_text(small_text),
                0.8,
                4.8,
                8.4,
                0.4,
                font_size=14,
                color=self.gray_50,
                align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
            )

    def _add_cards(self, slide, cards: list[Tag], top: float):
        """Add a row of cards to the slide."""
        num_cards = len(cards)
        if num_cards == 0:
            return

        card_width = 2.6
        gap = 0.3
        start_left = 0.8

        for i, card in enumerate(cards):
            title_el = card.find(class_="card-title")
            text_el = card.find(class_="card-text")
            number_el = card.find(class_="card-number")

            title = self.get_text(title_el)
            text = self.get_text(text_el)
            left = start_left + i * (card_width + gap)

            # Card background
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(card_width),
                Inches(1.8),
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = self.dark_gray
            card_shape.line.color.rgb = self.border_gray
            card_shape.line.width = Pt(1)

            if number_el:
                # Big number card
                number = self.get_text(number_el)
                self._add_number_card(slide, number, title, text, left, top, card_width)
            else:
                # Regular card with title and text
                self.add_text_box(
                    slide,
                    title,
                    left + 0.15,
                    top + 0.15,
                    card_width - 0.3,
                    0.4,
                    font_size=16,
                    bold=True,
                    color=self.primary,
                )
                self.add_text_box(
                    slide,
                    text,
                    left + 0.15,
                    top + 0.5,
                    card_width - 0.3,
                    1.2,
                    font_size=12,
                    color=self.gray_70,
                )

    def _add_number_card(
        self,
        slide,
        number: str,
        title: str,
        text: str,
        left: float,
        top: float,
        width: float,
    ):
        """Add a card with big number."""
        # Number
        self.add_text_box(
            slide,
            number,
            left + 0.1,
            top + 0.1,
            width - 0.2,
            0.7,
            font_size=48,
            bold=True,
            color=self.accent,
            align=PP_ALIGN.CENTER,
        )
        # Title
        self.add_text_box(
            slide,
            title,
            left + 0.1,
            top + 0.8,
            width - 0.2,
            0.3,
            font_size=14,
            bold=True,
            color=self.primary,
            align=PP_ALIGN.CENTER,
        )
        # Text
        self.add_text_box(
            slide,
            text,
            left + 0.1,
            top + 1.1,
            width - 0.2,
            0.5,
            font_size=10,
            color=self.gray_70,
            align=PP_ALIGN.CENTER,
        )

    def _add_tenets(self, slide, tenets: list[Tag], top: float):
        """Add tenet boxes to the slide."""
        num_tenets = len(tenets)

        if num_tenets >= 4:
            # Two columns
            for i, tenet in enumerate(tenets):
                col = i % 2
                row = i // 2
                self._add_single_tenet(
                    slide, tenet, 0.8 + col * 4.5, top + row * 1.0, width=4.2
                )
        else:
            # Single column
            for i, tenet in enumerate(tenets):
                self._add_single_tenet(slide, tenet, 0.8, top + i * 1.0, width=8.4)

    def _add_single_tenet(
        self, slide, tenet: Tag, left: float, top: float, width: float
    ):
        """Add a single tenet box with left border accent."""
        title = self.get_text(tenet.find(class_="tenet-title"))
        text = self.get_text(tenet.find(class_="tenet-text"))

        # Background box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.9)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x0D, 0x1A, 0x0D)
        box.line.fill.background()

        # Left accent bar
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.05), Inches(0.9)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.primary
        accent.line.fill.background()

        # Title
        self.add_text_box(
            slide,
            title,
            left + 0.15,
            top + 0.1,
            width - 0.3,
            0.3,
            font_size=14,
            bold=True,
            color=self.white,
        )

        # Text
        self.add_text_box(
            slide,
            text,
            left + 0.15,
            top + 0.4,
            width - 0.3,
            0.4,
            font_size=11,
            color=self.gray_70,
        )

    def _add_stats(self, slide, stat_grid: Tag, top: float):
        """Add statistics grid with big numbers."""
        stats = stat_grid.find_all(class_="stat")
        if not stats:
            return

        num_stats = len(stats)
        width_per_stat = 8.4 / num_stats
        start_left = 0.8

        for i, stat in enumerate(stats):
            number = self.get_text(stat.find(class_="stat-number"))
            label = self.get_text(stat.find(class_="stat-label"))
            left = start_left + i * width_per_stat

            # Big number
            self.add_text_box(
                slide,
                number,
                left,
                top,
                width_per_stat,
                0.6,
                font_size=40,
                bold=True,
                color=self.accent,
                align=PP_ALIGN.CENTER,
            )

            # Label
            self.add_text_box(
                slide,
                label,
                left,
                top + 0.6,
                width_per_stat,
                0.4,
                font_size=12,
                color=self.gray_70,
                align=PP_ALIGN.CENTER,
            )

    def _add_table(self, slide, table: Tag, top: float):
        """Add a data table to the slide."""
        rows = table.find_all("tr")
        if not rows:
            return

        row_height = 0.32
        for row_idx, row in enumerate(rows):
            cells = row.find_all(["th", "td"])
            is_header = row.find("th") is not None

            num_cols = len(cells)
            col_widths = [8.0 / num_cols] * num_cols
            if num_cols >= 3:
                col_widths = [2.5, 2.5, 3.0][:num_cols]

            left = 0.8
            for col_idx, cell in enumerate(cells):
                text = self.get_text(cell)
                width = col_widths[col_idx] if col_idx < len(col_widths) else 2.0

                if is_header:
                    color = self.primary
                    font_size = 12
                    bold = True
                else:
                    color = self.white if col_idx == 0 else self.gray_70
                    bold = col_idx == 0
                    font_size = 11

                self.add_text_box(
                    slide,
                    text,
                    left,
                    top + row_idx * row_height,
                    width,
                    row_height,
                    font_size=font_size,
                    bold=bold,
                    color=color,
                )
                left += width

    def _add_feature_list(self, slide, feature_list: Tag, top: float):
        """Add a feature list to the slide."""
        items = feature_list.find_all("li")
        for i, item in enumerate(items):
            text = self.get_text(item)

            # Determine color from content
            color = self.white
            if "✓" in text:
                color = RGBColor(0x00, 0xCC, 0x6A)  # Green
            elif "✗" in text:
                color = RGBColor(0xFF, 0x45, 0x3A)  # Red

            self.add_text_box(
                slide, text, 0.8, top + i * 0.4, 8.4, 0.4, font_size=16, color=color
            )

    def _add_generic_content(self, slide, slide_div: Tag, top: float, is_centered: bool):
        """
        Extract generic HTML content not using Amplifier Stories classes.
        This is a fallback for custom HTML that uses plain elements.
        """
        current_top = top
        
        # Track what we've already processed to avoid duplicates
        processed_elements = set()
        
        # Add the headline and section label elements we already processed
        headline = slide_div.find(["h1", "h2"], class_="headline") or slide_div.find("h1") or slide_div.find("h2")
        if headline:
            processed_elements.add(id(headline))
        
        section_label = slide_div.find(class_="section-label")
        if section_label:
            processed_elements.add(id(section_label))
        
        # Find all potentially relevant elements
        for element in slide_div.find_all(["h3", "ul", "div"]):
            # Skip if already processed
            if id(element) in processed_elements:
                continue
            
            # Skip if inside already-processed containers
            if (element.find_parent(class_=["card", "stat-grid", "stats-grid", "tenet", 
                                           "feature-list", "data-table", "versus", 
                                           "small-text", "image-container"])):
                continue
            
            # Handle h3 elements
            if element.name == "h3":
                text = self.get_text(element)
                if text:
                    self.add_text_box(
                        slide,
                        text,
                        0.8,
                        current_top,
                        8.4,
                        0.5,
                        font_size=28,
                        bold=True,
                        color=self.accent,
                        align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
                    )
                    current_top += 0.6
                    processed_elements.add(id(element))
            
            # Handle plain ul elements (not .feature-list)
            elif element.name == "ul" and "feature-list" not in element.get("class", []):
                items = element.find_all("li", recursive=False)
                if items:
                    for i, item in enumerate(items):
                        text = self.get_text(item)
                        if text:
                            # Use bullet points for generic lists
                            self.add_text_box(
                                slide,
                                "• " + text,
                                0.8,
                                current_top,
                                8.4,
                                0.3,
                                font_size=18,
                                color=self.white,
                            )
                            current_top += 0.35
                    processed_elements.add(id(element))
            
            # Handle divs with direct text content (like .flop-title, .flop-detail)
            elif element.name == "div":
                # Only process divs with classes that might contain content
                classes = element.get("class", [])
                # Skip structural containers
                if any(c in ["slide", "slides-container", "thirds", "halves", "fourths"] for c in classes):
                    continue
                
                # Get direct text (not from nested divs)
                direct_text = ""
                for child in element.children:
                    if isinstance(child, str):
                        direct_text += child.strip()
                
                if direct_text or (not list(element.find_all(["div", "ul", "table"])) and element.get_text(strip=True)):
                    text = self.get_text(element)
                    if text and len(text) > 3:  # Skip empty or trivial divs
                        # Check if it looks like a title (has title-like classes or is short and bold)
                        is_title = any(word in " ".join(classes).lower() for word in ["title", "heading", "header"])
                        
                        if is_title or len(text) < 60:
                            # Treat as a subheading
                            self.add_text_box(
                                slide,
                                text,
                                0.8,
                                current_top,
                                8.4,
                                0.4,
                                font_size=24,
                                bold=True,
                                color=self.accent,
                            )
                            current_top += 0.5
                        else:
                            # Treat as regular content
                            self.add_text_box(
                                slide,
                                text,
                                0.8,
                                current_top,
                                8.4,
                                0.3,
                                font_size=16,
                                color=self.gray_70,
                            )
                            current_top += 0.35
                        
                        processed_elements.add(id(element))

    def convert(self) -> Presentation:
        """Convert HTML to PowerPoint presentation."""
        slides = self.soup.find_all("div", class_="slide")

        if not slides:
            # Fallback: treat entire body as one slide
            body = self.soup.find("body")
            if body:
                slide = self.prs.slides.add_slide(self.blank_layout)
                self.set_background(slide)

                # Extract basic text content
                current_top = 0.6
                for element in body.find_all(["h1", "h2", "h3", "p", "ul"]):
                    if element.name in ["h1", "h2", "h3"]:
                        size = {"h1": 48, "h2": 36, "h3": 28}[element.name]
                        self.add_text_box(
                            slide,
                            self.get_text(element),
                            0.8,
                            current_top,
                            8.4,
                            0.6,
                            font_size=size,
                            bold=True,
                            color=self.white,
                        )
                        current_top += 0.7
                    elif element.name == "p":
                        self.add_text_box(
                            slide,
                            self.get_text(element),
                            0.8,
                            current_top,
                            8.4,
                            0.4,
                            font_size=18,
                            color=self.gray_70,
                        )
                        current_top += 0.5
                    elif element.name == "ul":
                        items = element.find_all("li")
                        for item in items:
                            self.add_text_box(
                                slide,
                                "• " + self.get_text(item),
                                0.8,
                                current_top,
                                8.4,
                                0.3,
                                font_size=16,
                                color=self.white,
                            )
                            current_top += 0.35
        else:
            # Process each .slide div
            for slide_div in slides:
                self.process_slide(slide_div)

        return self.prs


def main():
    parser = argparse.ArgumentParser(
        description="Convert HTML presentations to PowerPoint"
    )
    parser.add_argument("input", help="Input HTML file")
    parser.add_argument("output", help="Output PPTX file")
    parser.add_argument("--primary-color", help="Primary color (hex, no # prefix)")
    parser.add_argument("--accent-color", help="Accent color (hex, no # prefix)")

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_path = Path(args.output)

    print(f"Converting: {input_path}")
    print(f"Output: {output_path}")

    if args.primary_color:
        print(f"Primary color: #{args.primary_color}")
    if args.accent_color:
        print(f"Accent color: #{args.accent_color}")

    # Read HTML
    html_content = input_path.read_text(encoding="utf-8")

    # Convert
    converter = HTMLToPPTXConverter(
        html_content, primary_color=args.primary_color, accent_color=args.accent_color
    )
    prs = converter.convert()
    prs.save(str(output_path))

    print(f"✓ Done! Created {output_path}")


if __name__ == "__main__":
    main()
